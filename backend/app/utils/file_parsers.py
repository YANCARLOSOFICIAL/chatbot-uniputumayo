import logging
import re
from xml.etree import ElementTree as ET

logger = logging.getLogger(__name__)

# Map of file extension → canonical type key used internally
_EXT_MAP: dict[str, str] = {
    "pdf":  "pdf",
    "docx": "docx",
    "doc":  "docx",   # python-docx can open modern .doc; very old binary .doc will fail gracefully
    "txt":  "txt",
    "md":   "txt",
    "xlsx": "xlsx",
    "xls":  "xls",
    "csv":  "csv",
    "pptx": "pptx",
}

# Human-readable list for error messages
SUPPORTED_EXTENSIONS: list[str] = sorted(_EXT_MAP.keys())


def normalize_extension(ext: str) -> str | None:
    """Return canonical file type for a file extension, or None if unsupported."""
    return _EXT_MAP.get(ext.lower().lstrip("."))


def extract_text(file_path: str, file_type: str) -> str:
    """Dispatch to the appropriate parser for the given canonical file type."""
    dispatch = {
        "pdf":  _extract_pdf,
        "docx": _extract_docx,
        "txt":  _extract_txt,
        "xlsx": _extract_xlsx,
        "xls":  _extract_xls,
        "csv":  _extract_csv,
        "pptx": _extract_pptx,
    }
    fn = dispatch.get(file_type)
    if fn is None:
        raise ValueError(f"Tipo de archivo no soportado: '{file_type}'")
    return fn(file_path)


# ── PDF ─────────────────────────────────────────────────────────────────────

def _extract_pdf(file_path: str) -> str:
    """Coordinate-based extraction that reconstructs reading order for complex layouts.

    Standard PyMuPDF get_text("text") scrambles columns and rotated grids (e.g.
    curriculum pensum tables).  This approach rebuilds rows from (x, y) spans.
    """
    import fitz  # PyMuPDF

    all_pages = []
    with fitz.open(file_path) as doc:
        for page in doc:
            page_dict = page.get_text("dict")
            words: list[tuple[float, float, str]] = []

            for block in page_dict.get("blocks", []):
                if block.get("type") != 0:  # skip image blocks
                    continue
                for line in block.get("lines", []):
                    for span in line.get("spans", []):
                        text = span.get("text", "").strip()
                        if not text:
                            continue
                        if len(text) == 1 and not text.isalpha():
                            continue
                        words.append((span["bbox"][0], span["bbox"][1], text))

            if not words:
                continue

            ROW_BUCKET = 6
            rows: dict[int, list[tuple[float, str]]] = {}
            for x0, y0, text in words:
                key = round(y0 / ROW_BUCKET) * ROW_BUCKET
                rows.setdefault(key, []).append((x0, text))

            page_lines = []
            for y_key in sorted(rows):
                row_words = sorted(rows[y_key], key=lambda w: w[0])
                line = " ".join(w for _, w in row_words).strip()
                meaningful = sum(1 for c in line if c.isalpha())
                if len(line) > 3 and meaningful >= 2:
                    page_lines.append(line)

            if page_lines:
                all_pages.append("\n".join(page_lines))

    return "\n\n".join(all_pages)


# ── DOCX ─────────────────────────────────────────────────────────────────────

def _extract_docx(file_path: str) -> str:
    from docx import Document

    doc = Document(file_path)
    parts = []

    for paragraph in doc.paragraphs:
        if paragraph.text.strip():
            parts.append(paragraph.text)

    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            non_empty = [c for c in cells if c]
            # >= 1 — a merged/spanning header cell in a docx table has the
            # same single-populated-cell shape as the xlsx case; see
            # _extract_xlsx. This one was missed when that fix was applied
            # everywhere else (xls, csv, pptx already use >= 1).
            if len(non_empty) >= 1:
                parts.append(" | ".join(non_empty))

    return "\n\n".join(parts)


# ── TXT / Markdown ────────────────────────────────────────────────────────────

def _extract_txt(file_path: str) -> str:
    """Try multiple encodings; Spanish files from Windows are often cp1252/latin-1."""
    for enc in ("utf-8-sig", "utf-8", "latin-1", "cp1252"):
        try:
            with open(file_path, "r", encoding=enc) as f:
                return f.read()
        except (UnicodeDecodeError, LookupError):
            continue
    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        logger.warning("TXT file %s decoded with replacement characters", file_path)
        return f.read()


# ── Excel XLSX ────────────────────────────────────────────────────────────────

def _xlsx_cell_str(val) -> str:
    """Convert an openpyxl cell value to a clean string."""
    if val is None:
        return ""
    if isinstance(val, float) and val == int(val):
        return str(int(val))
    return str(val).strip()


def _looks_numeric(s: str) -> bool:
    try:
        float(s)
        return True
    except ValueError:
        return False


_ROMAN_RE = re.compile(r"^[IVXLCDM]+$")

# Roman numeral → word ordinal, for semester-grid labels. Curriculum grids
# label semesters "I".."X" as textbox shapes (see _semester_column_ranges),
# but a real question is never phrased that way ("segundo semestre", never
# "semestre II") — confirmed live: retrieval for "segundo semestre" pulled in
# semesters I and IV instead of II, purely because nothing in the chunk text
# shared any token with the query's ordinal. Embedding cosine similarity and
# the keyword-overlap rerank both need a literal match to work with.
_ROMAN_TO_ORDINAL = {
    "I": "primer", "II": "segundo", "III": "tercer", "IV": "cuarto",
    "V": "quinto", "VI": "sexto", "VII": "séptimo", "VIII": "octavo",
    "IX": "noveno", "X": "décimo", "XI": "undécimo", "XII": "duodécimo",
}

_NS_MAIN = "http://schemas.openxmlformats.org/spreadsheetml/2006/main"
_NS_DOC_REL = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
_NS_PKG_REL = "http://schemas.openxmlformats.org/package/2006/relationships"
_NS_XDR = "http://schemas.openxmlformats.org/drawingml/2006/spreadsheetDrawing"
_NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"


def _sheet_drawing_parts(zf) -> dict[str, str]:
    """Map each sheet's visible name to its drawing XML part, if it has one.

    Needed because semester labels on a "malla curricular" grid (see
    `_semester_column_ranges`) are floating textbox shapes, not cell values —
    openpyxl's cell API can't see them at all, only the raw drawing XML can.
    """
    wb_xml = ET.fromstring(zf.read("xl/workbook.xml"))
    wb_rels = ET.fromstring(zf.read("xl/_rels/workbook.xml.rels"))

    rid_to_target = {
        rel.get("Id"): rel.get("Target")
        for rel in wb_rels.findall(f"{{{_NS_PKG_REL}}}Relationship")
    }

    result: dict[str, str] = {}
    for sheet_el in wb_xml.findall(f".//{{{_NS_MAIN}}}sheet"):
        rid = sheet_el.get(f"{{{_NS_DOC_REL}}}id")
        target = rid_to_target.get(rid)
        name = sheet_el.get("name")
        if not target or not name:
            continue
        sheet_part = target if target.startswith("xl/") else f"xl/{target}"
        rels_part = sheet_part.rsplit("/", 1)[0] + "/_rels/" + sheet_part.rsplit("/", 1)[1] + ".rels"
        try:
            sheet_rels = ET.fromstring(zf.read(rels_part))
        except KeyError:
            continue
        for rel in sheet_rels.findall(f"{{{_NS_PKG_REL}}}Relationship"):
            if rel.get("Type", "").endswith("/drawing"):
                drawing_file = rel.get("Target").rsplit("/", 1)[-1]
                result[name] = f"xl/drawings/{drawing_file}"
                break

    return result


def _semester_column_ranges(zf, drawing_part: str) -> list[tuple[int, int, str]]:
    """Extract 1-indexed (start_col, end_col, label) ranges for Roman-numeral
    textbox shapes overlaid on a sheet (semester headers on a curriculum grid
    where each semester is a block of columns, not a row) — see docstring on
    `_extract_xlsx` for why this exists at all.
    """
    try:
        root = ET.fromstring(zf.read(drawing_part))
    except KeyError:
        return []

    ranges: list[tuple[int, int, str]] = []
    for anchor in root.findall(f"{{{_NS_XDR}}}twoCellAnchor"):
        frm, to = anchor.find(f"{{{_NS_XDR}}}from"), anchor.find(f"{{{_NS_XDR}}}to")
        text_el = anchor.find(f".//{{{_NS_A}}}t")
        if frm is None or to is None or text_el is None or not text_el.text:
            continue
        label = text_el.text.strip()
        if not _ROMAN_RE.match(label):
            continue  # not a semester label — an arrow, box, or other decoration
        from_col_el, to_col_el = frm.find(f"{{{_NS_XDR}}}col"), to.find(f"{{{_NS_XDR}}}col")
        if from_col_el is None or to_col_el is None:
            continue
        ranges.append((int(from_col_el.text) + 1, int(to_col_el.text) + 1, label))

    ranges.sort(key=lambda r: r[0])
    # A single stray Roman-numeral shape elsewhere on an otherwise normal sheet
    # isn't a semester grid — need at least 2 blocks for the column-range
    # split below to mean anything.
    return ranges if len(ranges) >= 2 else []


def _semester_column_ranges_from_cells(
    ws,
) -> tuple[list[tuple[int, int, str]], int] | None:
    """Same idea as `_semester_column_ranges`, but for mallas where the
    semester header is a literal cell value (e.g. "I" typed into a cell)
    rather than a floating textbox shape overlaid on the sheet.

    Both styles are real, common institutional exports — confirmed live on
    two different programs' curriculum files (Contaduría, Gastronomía) using
    plain cell values instead of shapes. Without this fallback those sheets
    have no semester-grid detected at all and fall through to the plain
    row-based path below, which — for a "semester is a column block" layout
    — interleaves every semester's courses onto the same line (the exact
    failure `_semester_column_ranges` exists to prevent for the shape-based
    style).

    Also returns the 1-indexed row the header was found on: front-matter
    above it (title blocks, the table's own column-header row) sits in the
    same left-margin columns real per-section category labels use further
    down, and would otherwise be picked up as if it were one — confirmed
    live on a real file where a column header ("CREDITOS POR CICLO") ended
    up glued onto every single course's category label.
    """
    for row_idx, row in enumerate(ws.iter_rows(), start=1):
        matches = [
            (cell.column, str(cell.value).strip())
            for cell in row
            if cell.value is not None and _ROMAN_RE.match(str(cell.value).strip())
        ]
        if len(matches) < 2:
            continue
        ranges: list[tuple[int, int, str]] = []
        for i, (col, label) in enumerate(matches):
            end_col = matches[i + 1][0] - 1 if i + 1 < len(matches) else ws.max_column
            ranges.append((col, end_col, label))
        return ranges, row_idx
    return None


def _load_semester_shape_map(file_path: str) -> dict[str, list[tuple[int, int, str]]]:
    import zipfile

    try:
        with zipfile.ZipFile(file_path) as zf:
            drawing_parts = _sheet_drawing_parts(zf)
            return {
                name: _semester_column_ranges(zf, part)
                for name, part in drawing_parts.items()
            }
    except Exception as e:
        # Best-effort enhancement only — any surprise here (different Excel
        # generator, no drawings, unexpected XML shape) just means every
        # sheet falls back to the plain row-based extraction below.
        logger.debug("Could not read semester-grid shape labels (%s): %s", file_path, e)
        return {}


def _extract_semester_grid_sheet(
    ws,
    sheet_name: str,
    semester_ranges: list[tuple[int, int, str]],
    category_start_row: int = 1,
) -> list[str]:
    """Extract a "malla curricular" sheet where each semester is a block of
    COLUMNS (course code/name/hours stacked vertically within it) rather than
    a row — plain row-by-row extraction concatenates every semester's courses
    onto the same line with no way to tell them apart (confirmed live: a RAG
    query for "segundo semestre" returned zero matches because the retrieved
    chunks were an unreadable mix of 6+ semesters' course names on one line).
    Splitting by each semester's own column range keeps them separate, the
    same way a normal one-row-per-course sheet already reads cleanly.
    """
    max_row = ws.max_row
    first_semester_col = semester_ranges[0][0]

    # The grid ends with a "TOTAL" row (per-semester credit sums) followed by
    # a "REQUISITOS DE GRADO" legend/footnotes block. Both are real
    # information, but neither can be handled like a normal course row:
    # - The legend text is identical no matter which semester's columns it's
    #   sliced from, so treating it like a course row would repeat it
    #   verbatim in every semester's chunk — confirmed live to actively break
    #   retrieval (a query for "materias de formación básica" returned only
    #   this repeated footer, no real courses, because it out-scored every
    #   semester's actual content). It's pulled into its own section instead,
    #   appearing once.
    # - The TOTAL row genuinely IS per-semester data (each semester's own
    #   credit sum lives in that semester's columns) — but as a *separate*
    #   section it's one long, mostly-numeric row spanning all semesters,
    #   which (confirmed live) still out-scores real course content for
    #   broader questions ("cuántos créditos tiene el programa"). It reads
    #   correctly, and stays out of competition with course chunks, once
    #   folded into the semester section it actually belongs to instead.
    # The "TOTAL"/"REQUISITOS DE GRADO" label doesn't always sit in column 1 —
    # confirmed live on a real file where it's in column 2 instead (with
    # column 1 empty on that row). Scanning every column left of the first
    # semester block (the same range category_by_row already reads) instead
    # of hardcoding column 1 catches both layouts.
    total_row: int | None = None
    footer_start_row = max_row + 1
    for r in range(1, max_row + 1):
        row_labels = [
            _xlsx_cell_str(ws.cell(row=r, column=c).value).upper()
            for c in range(1, first_semester_col)
        ]
        # Exact "TOTAL" works for a plain total row, but not every file's
        # total row is that bare — confirmed live on one that reads "TOTAL
        # CREDITOS CICLO TECNOLOGICO" instead. A prefix check catches both
        # without matching real category text (none of it starts with the
        # word "TOTAL" in any file seen so far).
        is_total_row = any(lbl.startswith("TOTAL") for lbl in row_labels)
        if is_total_row:
            total_row = r
        if is_total_row or any("REQUISITOS DE GRADO" in lbl for lbl in row_labels):
            footer_start_row = r
            break

    course_rows_end = total_row if total_row is not None else footer_start_row

    # Row-category labels (e.g. "FORMACIÓN BÁSICA") live in the columns to the
    # left of the first semester block, in cells merged down across many rows
    # — carry the last-seen value forward per row, same as a merged cell reads
    # in every other row-based extractor in this file. Keyed by the actual
    # column number, not by scan position — a sibling category (e.g.
    # "COMPONENTE BÁSICO" replaced later by "COMPONENTE PROFESIONAL", both in
    # the same column) must overwrite its own slot even when an earlier,
    # unrelated column (e.g. the leftmost one) is never populated at all.
    # Keying by scan position instead let the old value keep accumulating
    # alongside the new one instead of being replaced — confirmed live on a
    # real file where column 1 is always blank: rows ended up labeled
    # "COMPONENTE BÁSICO - COMPONENTE PROFESIONAL - COMPONENTE TÉCNICO
    # PRODUCCIÓN" all at once, mixing categories that are mutually exclusive.
    category_by_row: dict[int, str] = {}
    running: dict[int, str] = {}
    for r in range(category_start_row, course_rows_end):
        for c in range(1, first_semester_col):
            raw = ws.cell(row=r, column=c).value
            val = _xlsx_cell_str(raw)
            # Some layouts pack a per-category credit count / percentage
            # (e.g. "28", "0.264...") into these same left-margin columns,
            # alongside the actual text labels — confirmed live on a real
            # file where that produced a "category" like "AREA DE FORMACION
            # BASICA - ... - 28 - 0.264...". A bare number here is summary
            # data, never a category name, so it's excluded from the label.
            if val and not isinstance(raw, (int, float)):
                running[c] = val
        category_by_row[r] = " - ".join(running[c] for c in sorted(running))

    parts: list[str] = []
    for start_col, end_col, label in semester_ranges:
        # Credits/hours in this grid are always whole numbers, so a column
        # that ever holds a fractional value anywhere in the course-row range
        # isn't part of the course grid at all — it's a leaked calculation
        # artifact. Confirmed live: a hidden helper column (used to compute
        # the sheet's own %-of-program breakdown, e.g. "14.516129032258064",
        # alongside its own whole-number counterpart like "9") landed inside
        # the last semester's detected column range purely by position,
        # contaminating that semester's chunk with unrelated figures.
        # Dropping the whole column (not just the fractional rows) once any
        # fractional value is seen in it also clears the paired whole-number
        # rows, which look exactly like a real credit value in isolation.
        tainted_cols = {
            c for c in range(start_col, end_col + 1)
            if any(
                isinstance(v := ws.cell(row=r, column=c).value, float) and v != int(v)
                for r in range(1, course_rows_end)
            )
        }

        rows_text: list[str] = []
        semester_credits = 0
        for r in range(1, course_rows_end):
            cells = [
                "" if c in tainted_cols else _xlsx_cell_str(ws.cell(row=r, column=c).value)
                for c in range(start_col, end_col + 1)
            ]
            if not any(cells):
                continue
            non_empty = [c for c in cells if c]
            if non_empty == [label]:
                # The semester's own header row, when it's a literal cell
                # value rather than a floating shape (see
                # _semester_column_ranges_from_cells) — it's already shown in
                # this section's own header line, so keeping it here too
                # would just repeat it as a floating, context-less row.
                continue
            # A course's "numbers" row always carries its credit value, but
            # not always in the same shape — confirmed live across three
            # different real templates: "4 | 8 | T | 4" (type marker
            # second-to-last), "3 |  | 3" (no marker, just numbers), and
            # "HP | 3 | HTI | 3 | TP | 3 | CR | 3" (label/value pairs, credit
            # is the value right after the literal "CR" label). Checking for
            # a literal "CR" label first catches that third shape directly;
            # falling back to the positional heuristics below covers the
            # other two instead of silently summing to 0 for whichever shape
            # isn't handled. Summing these directly from the course rows
            # already being read here is self-consistent with what the
            # chunk itself says, unlike trusting the sheet's own TOTAL row:
            # that row's columns don't line up 1:1 with each semester's
            # course columns, and naively taking the first non-blank cell in
            # it overcounted every semester (confirmed live: 15 instead of
            # the real 13 for Semestre I, verified by hand-summing that
            # semester's courses).
            credit_str: str | None = None
            if "CR" in non_empty and non_empty.index("CR") + 1 < len(non_empty):
                credit_str = non_empty[non_empty.index("CR") + 1]
            elif len(non_empty) >= 2 and non_empty[-2] in ("T", "TP", "P"):
                credit_str = non_empty[-1]
            elif non_empty and all(_looks_numeric(v) for v in non_empty):
                credit_str = non_empty[-1]
            if credit_str is not None:
                try:
                    semester_credits += int(float(credit_str))
                except ValueError:
                    pass
            line = " | ".join(cells).strip(" |")
            category = category_by_row.get(r, "")
            rows_text.append(f"{category} | {line}" if category else line)

        if semester_credits:
            rows_text.append(f"Total créditos de este semestre: {semester_credits}")

        if rows_text:
            ordinal = _ROMAN_TO_ORDINAL.get(label)
            header_label = f"SEMESTRE {label} ({ordinal} semestre)" if ordinal else f"SEMESTRE {label}"
            parts.append(f"=== HOJA: {sheet_name} — {header_label} ===")
            parts.append("\n".join(rows_text))

    if footer_start_row <= max_row:
        # The footer isn't one flowing block of text — it's several unrelated
        # tables placed side by side in disjoint column ranges that just
        # happen to share the same rows (an abbreviation legend, per-cycle
        # graduation requirements, the electives portfolio, program-wide
        # credit totals). Joining every column of every row as if it were a
        # single table interleaves them into noise — confirmed live: rows
        # from the electives list came out mixed with unrelated legend cells
        # ("CODIG | | PR. | CODIG | ... | 2 | Diseño de Apps"). Grouping
        # columns into contiguous runs (a run ends at the first fully-empty
        # column) and keeping each run's rows in their own section stops the
        # unrelated tables from bleeding into each other.
        footer_row_range = range(footer_start_row, max_row + 1)
        occupied_cols = [
            c for c in range(1, ws.max_column + 1)
            if any(
                _xlsx_cell_str(ws.cell(row=r, column=c).value)
                for r in footer_row_range if r != total_row
            )
        ]
        col_blocks: list[list[int]] = []
        for c in occupied_cols:
            if col_blocks and c == col_blocks[-1][-1] + 1:
                col_blocks[-1].append(c)
            else:
                col_blocks.append([c])

        for block in col_blocks:
            block_rows: list[str] = []
            for r in footer_row_range:
                if r == total_row:
                    continue
                cells = [_xlsx_cell_str(ws.cell(row=r, column=c).value) for c in block]
                if any(cells):
                    block_rows.append(" | ".join(cells).strip(" |"))
            # A block with no label at all — just a bare number, e.g. a
            # stray per-semester sub-total column whose own header column
            # fell outside this contiguous run — carries no retrievable
            # meaning on its own, so it's dropped instead of surfacing as a
            # floating, context-less digit.
            if block_rows and any(re.search(r"[^\W\d_]", row) for row in block_rows):
                parts.append(f"=== HOJA: {sheet_name} — INFORMACIÓN COMPLEMENTARIA ===")
                parts.append("\n".join(block_rows))

    return parts


def _extract_xlsx(file_path: str) -> str:
    """Extract all sheets from an XLSX file as labeled pipe-separated rows.

    data_only=True returns computed values instead of formulas, which is what we
    want for embedding (e.g. curricula with calculated credit totals).

    Most institutional spreadsheets are simple one-row-per-record tables, where
    reading top-to-bottom/left-to-right and pipe-joining each row already
    produces something an embedding model and an LLM can use. But some
    curriculum files ("mallas") lay semesters out as column BLOCKS instead —
    course code/name/hours stacked vertically within each semester's columns,
    with the semester number itself drawn as a floating textbox shape rather
    than a cell value. Flat row extraction on one of these interleaves every
    semester's courses onto the same line with nothing to tell them apart.
    `_load_semester_shape_map` detects this shape (via the sheet's own
    embedded drawing XML) and, when found, `_extract_semester_grid_sheet`
    splits that sheet by semester column-range instead of by row.
    """
    import openpyxl

    wb = openpyxl.load_workbook(file_path, data_only=True)
    semester_map = _load_semester_shape_map(file_path)
    parts = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        semester_ranges = semester_map.get(sheet_name)
        category_start_row = 1
        if not semester_ranges:
            detected = _semester_column_ranges_from_cells(ws)
            if detected:
                semester_ranges, category_start_row = detected

        if semester_ranges:
            parts.extend(
                _extract_semester_grid_sheet(ws, sheet_name, semester_ranges, category_start_row)
            )
            continue

        rows_text: list[str] = []
        for row in ws.iter_rows(values_only=True):
            cells = [_xlsx_cell_str(c) for c in row]
            non_empty = [c for c in cells if c]
            # >= 1, not >= 2: a merged cell (e.g. a "SEMESTRE I" header spanning
            # several columns) has openpyxl return that value in exactly one cell
            # and None in the rest of the merged range — a >= 2 threshold would
            # silently drop every such header row, losing the section structure
            # a curriculum spreadsheet relies on.
            if len(non_empty) >= 1:
                rows_text.append(" | ".join(cells).strip(" |"))

        if rows_text:
            parts.append(f"=== HOJA: {sheet_name} ===")
            parts.append("\n".join(rows_text))

    return "\n\n".join(parts)


# ── Excel XLS (legacy 97-2003) ────────────────────────────────────────────────

class _XlrdCellAdapter:
    """A single cell, adapted to expose `.value` (converted the same way as
    the plain xlrd extraction path below) and a 1-indexed `.column`.
    """

    def __init__(self, cell, datemode, column: int):
        self._cell = cell
        self._datemode = datemode
        self.column = column

    @property
    def value(self):
        import xlrd

        cell = self._cell
        if cell.ctype == xlrd.XL_CELL_EMPTY or cell.value == "":
            return None
        if cell.ctype == xlrd.XL_CELL_DATE:
            try:
                return xlrd.xldate_as_datetime(cell.value, self._datemode).strftime("%Y-%m-%d")
            except Exception:
                return str(cell.value)
        if cell.ctype == xlrd.XL_CELL_NUMBER:
            v = cell.value
            return int(v) if v == int(v) else v
        return cell.value


class _XlrdSheetAdapter:
    """Adapts an xlrd sheet to the 1-indexed `.cell(row=, column=)` /
    `.max_row` / `.max_column` / `.iter_rows()` interface openpyxl exposes,
    so `_extract_semester_grid_sheet` and `_semester_column_ranges_from_cells`
    — already battle-tested on real xlsx curriculum files — work unmodified
    on legacy .xls files too, instead of duplicating that logic for a
    second cell-access API.
    """

    def __init__(self, sheet, datemode):
        self._sheet = sheet
        self._datemode = datemode

    @property
    def max_row(self) -> int:
        return self._sheet.nrows

    @property
    def max_column(self) -> int:
        return self._sheet.ncols

    def cell(self, row: int, column: int) -> _XlrdCellAdapter:
        return _XlrdCellAdapter(self._sheet.cell(row - 1, column - 1), self._datemode, column)

    def iter_rows(self):
        for r in range(self._sheet.nrows):
            yield [
                _XlrdCellAdapter(self._sheet.cell(r, c), self._datemode, c + 1)
                for c in range(self._sheet.ncols)
            ]


def _extract_xls(file_path: str) -> str:
    """Extract an old-format .xls workbook using xlrd."""
    import xlrd

    wb = xlrd.open_workbook(file_path)
    parts = []

    for sheet_idx in range(wb.nsheets):
        ws = wb.sheet_by_index(sheet_idx)

        adapter = _XlrdSheetAdapter(ws, wb.datemode)
        detected = _semester_column_ranges_from_cells(adapter)
        if detected:
            semester_ranges, category_start_row = detected
            parts.extend(
                _extract_semester_grid_sheet(adapter, ws.name, semester_ranges, category_start_row)
            )
            continue

        rows_text: list[str] = []

        for r in range(ws.nrows):
            cells: list[str] = []
            for c_idx in range(ws.ncols):
                cell = ws.cell(r, c_idx)
                if cell.ctype == xlrd.XL_CELL_DATE:
                    try:
                        import datetime as _dt
                        val = xlrd.xldate_as_datetime(cell.value, wb.datemode).strftime("%Y-%m-%d")
                    except Exception:
                        val = str(cell.value)
                elif cell.ctype == xlrd.XL_CELL_NUMBER:
                    v = cell.value
                    val = str(int(v)) if v == int(v) else str(v)
                else:
                    val = str(cell.value).strip() if cell.value else ""
                cells.append(val)

            non_empty = [c for c in cells if c]
            # See the matching comment in _extract_xlsx — merged-cell section
            # headers (xlrd reports the same value-in-anchor-cell-only shape).
            if len(non_empty) >= 1:
                rows_text.append(" | ".join(cells).strip(" |"))

        if rows_text:
            parts.append(f"=== HOJA: {ws.name} ===")
            parts.append("\n".join(rows_text))

    return "\n\n".join(parts)


# ── CSV ───────────────────────────────────────────────────────────────────────

def _extract_csv(file_path: str) -> str:
    """Parse CSV with automatic encoding detection (covers UTF-8, BOM, Latin-1, cp1252)."""
    import csv

    encodings = ("utf-8-sig", "utf-8", "latin-1", "cp1252")
    for enc in encodings:
        try:
            rows_text: list[str] = []
            with open(file_path, newline="", encoding=enc) as f:
                reader = csv.reader(f)
                for row in reader:
                    cells = [c.strip() for c in row]
                    non_empty = [c for c in cells if c]
                    # >= 1 — a CSV exported from a spreadsheet with merged header
                    # cells (e.g. "SEMESTRE I,,,") has the same single-populated-
                    # column shape as the xlsx/xls case; see _extract_xlsx.
                    if len(non_empty) >= 1:
                        rows_text.append(" | ".join(cells).strip(" |"))
            return "\n".join(rows_text)
        except (UnicodeDecodeError, csv.Error):
            continue

    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        return f.read()


# ── PowerPoint PPTX ───────────────────────────────────────────────────────────

def _extract_pptx(file_path: str) -> str:
    """Extract text from all slides, including shapes and tables."""
    from pptx import Presentation

    prs = Presentation(file_path)
    parts = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_texts: list[str] = []

        for shape in slide.shapes:
            # Text frames (titles, body, text boxes)
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    slide_texts.append(text)
            # Tables in slides
            if shape.has_table:
                for trow in shape.table.rows:
                    cells = [c.text.strip() for c in trow.cells]
                    non_empty = [c for c in cells if c]
                    # >= 1 — a merged/spanning header cell in a slide table has
                    # the same single-populated-cell shape as the xlsx case.
                    if len(non_empty) >= 1:
                        slide_texts.append(" | ".join(non_empty))

        if slide_texts:
            # "===" (not "---") — clean_text() strips runs of 3+ hyphens as
            # decorative separators, which would destroy a "---"-style marker
            # before chunk_tabular_text ever sees it.
            parts.append(f"=== DIAPOSITIVA {slide_num} ===")
            parts.append("\n".join(slide_texts))

    return "\n\n".join(parts)
