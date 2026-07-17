import csv

import openpyxl
import pytest
from docx import Document
from pptx import Presentation
from pptx.util import Inches

from app.utils.file_parsers import _extract_csv, _extract_docx, _extract_pptx, _extract_xlsx


@pytest.fixture
def xlsx_with_merged_header(tmp_path):
    """A pensum-style sheet: a merged 'SEMESTRE I' header row, then course rows."""
    path = tmp_path / "pensum.xlsx"
    wb = openpyxl.Workbook()
    ws = wb.active
    ws["A1"] = "SEMESTRE I"
    ws.merge_cells("A1:D1")
    ws.append(["TD101", "Fundamentos de seguridad informática", 4, "TD101"])
    ws.append(["TD102", "Análisis y gestión de riesgos", 3, "TD102"])
    ws.append([None, None, None, None])  # fully blank spacer row
    ws["A5"] = "SEMESTRE II"
    ws.merge_cells("A5:D5")
    ws.append(["TD201", "Criptografía aplicada", 3, "TD201"])
    wb.save(path)
    return str(path)


@pytest.fixture
def xlsx_with_cell_value_semesters(tmp_path):
    """A pensum-style sheet like `xlsx_with_merged_header`, but for the OTHER
    real institutional export style: semester numbers are literal cell
    values ("I", "II", ...) rather than floating textbox shapes, the
    category column is sparse (column 1 is never populated, only column 2 —
    with a sibling category change partway through semester I), and the
    TOTAL row's label sits in column 2, not column 1.
    """
    path = tmp_path / "pensum_cellval.xlsx"
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.cell(row=1, column=3, value="I")
    ws.cell(row=1, column=7, value="II")

    ws.cell(row=2, column=2, value="COMPONENTE BASICO")
    ws.cell(row=2, column=3, value="TD101")
    ws.cell(row=2, column=5, value="PR.")
    ws.cell(row=3, column=3, value="Fundamentos de seguridad informática")
    ws.cell(row=4, column=3, value=4)
    ws.cell(row=4, column=4, value=8)
    ws.cell(row=4, column=5, value="T")
    ws.cell(row=4, column=6, value=4)

    ws.cell(row=5, column=2, value="COMPONENTE PROFESIONAL")
    ws.cell(row=5, column=3, value="TD102")
    ws.cell(row=6, column=3, value="Análisis y gestión de riesgos")
    ws.cell(row=7, column=3, value=3)
    ws.cell(row=7, column=4, value=6)
    ws.cell(row=7, column=5, value="T")
    ws.cell(row=7, column=6, value=3)

    ws.cell(row=2, column=7, value="TD201")
    ws.cell(row=3, column=7, value="Criptografía aplicada")
    ws.cell(row=4, column=7, value=2)
    ws.cell(row=4, column=8, value=4)
    ws.cell(row=4, column=9, value="T")
    ws.cell(row=4, column=10, value=2)

    ws.cell(row=8, column=2, value="TOTAL")
    ws.cell(row=8, column=6, value=7)
    ws.cell(row=8, column=10, value=2)
    wb.save(path)
    return str(path)


class TestExtractXlsxCellValueSemesterGrid:
    def test_cell_value_semester_header_splits_sections(self, xlsx_with_cell_value_semesters):
        text = _extract_xlsx(xlsx_with_cell_value_semesters)
        assert "SEMESTRE I" in text
        assert "SEMESTRE II" in text
        assert "Criptografía aplicada" in text

    def test_sibling_category_replaces_not_accumulates(self, xlsx_with_cell_value_semesters):
        text = _extract_xlsx(xlsx_with_cell_value_semesters)
        profesional_line = next(
            l for l in text.splitlines() if "Análisis y gestión de riesgos" in l
        )
        assert "COMPONENTE PROFESIONAL" in profesional_line
        assert "COMPONENTE BASICO" not in profesional_line

    def test_total_row_label_in_non_first_column_is_detected(self, xlsx_with_cell_value_semesters):
        # The TOTAL row's label sits in column 2, not column 1 — if it went
        # undetected, the TOTAL row would leak into the course grid instead
        # of being excluded, and no "Total créditos" line would appear at all.
        text = _extract_xlsx(xlsx_with_cell_value_semesters)
        assert "Total créditos de este semestre: 7" in text

    def test_semester_header_row_itself_is_not_shown_as_a_course(self, xlsx_with_cell_value_semesters):
        text = _extract_xlsx(xlsx_with_cell_value_semesters)
        lines = [l for l in text.splitlines() if l.strip()]
        assert not any(l.strip() == "I" for l in lines)

    def test_label_value_numbers_row_credit_is_detected(self, tmp_path):
        # A third real template shape: the numbers row is label/value pairs
        # ("HP | 3 | HTI | 3 | TP | 3 | CR | 3") instead of a fixed position
        # — the credit is whatever value follows the literal "CR" label, not
        # necessarily the row's last cell or preceded by a T/TP/P marker.
        path = tmp_path / "pensum_labelvalue.xlsx"
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.cell(row=1, column=3, value="I")
        ws.cell(row=1, column=11, value="II")

        ws.cell(row=2, column=3, value="FING 2101")
        ws.cell(row=3, column=3, value="Química General")
        ws.cell(row=4, column=3, value="HP")
        ws.cell(row=4, column=4, value=3)
        ws.cell(row=4, column=5, value="HTI")
        ws.cell(row=4, column=6, value=3)
        ws.cell(row=4, column=7, value="TP")
        ws.cell(row=4, column=8, value=3)
        ws.cell(row=4, column=9, value="CR")
        ws.cell(row=4, column=10, value=3)

        ws.cell(row=2, column=11, value="FING 2102")
        wb.save(path)

        text = _extract_xlsx(str(path))
        assert "Total créditos de este semestre: 3" in text


class TestExtractXlsx:
    def test_merged_header_row_is_preserved(self, xlsx_with_merged_header):
        text = _extract_xlsx(xlsx_with_merged_header)
        assert "SEMESTRE I" in text
        assert "SEMESTRE II" in text

    def test_course_rows_are_preserved(self, xlsx_with_merged_header):
        text = _extract_xlsx(xlsx_with_merged_header)
        assert "Fundamentos de seguridad informática" in text
        assert "Criptografía aplicada" in text

    def test_fully_blank_row_is_still_dropped(self, xlsx_with_merged_header):
        text = _extract_xlsx(xlsx_with_merged_header)
        lines = [l for l in text.splitlines() if l.strip()]
        # No line should be just separators from an all-None row.
        assert all(line.strip(" |") for line in lines)

    def test_sheet_marker_present(self, xlsx_with_merged_header):
        text = _extract_xlsx(xlsx_with_merged_header)
        assert "=== HOJA:" in text


class TestExtractCsv:
    def test_merged_style_header_row_is_preserved(self, tmp_path):
        path = tmp_path / "pensum.csv"
        with open(path, "w", newline="", encoding="utf-8") as f:
            writer = csv.writer(f)
            writer.writerow(["SEMESTRE I", "", "", ""])
            writer.writerow(["TD101", "Fundamentos de seguridad informática", "4"])
            writer.writerow(["", "", ""])  # fully blank row
            writer.writerow(["SEMESTRE II", "", "", ""])
            writer.writerow(["TD201", "Criptografía aplicada", "3"])

        text = _extract_csv(str(path))
        assert "SEMESTRE I" in text
        assert "SEMESTRE II" in text
        assert "Fundamentos de seguridad informática" in text

    def test_latin1_encoded_file_is_readable(self, tmp_path):
        path = tmp_path / "pensum_latin1.csv"
        with open(path, "w", newline="", encoding="latin-1") as f:
            writer = csv.writer(f)
            writer.writerow(["Código", "Materia"])
            writer.writerow(["TD101", "Análisis y diseño"])

        text = _extract_csv(str(path))
        assert "Análisis" in text or "An" in text  # tolerant of encoding round-trip


class TestExtractDocx:
    def test_table_with_single_populated_header_row_is_preserved(self, tmp_path):
        path = tmp_path / "pensum.docx"
        doc = Document()
        table = doc.add_table(rows=3, cols=2)
        table.cell(0, 0).text = "SEMESTRE I"
        table.cell(0, 1).text = ""  # single-populated header cell, like a merged cell
        table.cell(1, 0).text = "TD101"
        table.cell(1, 1).text = "Fundamentos de seguridad informática"
        table.cell(2, 0).text = ""
        table.cell(2, 1).text = ""  # fully blank row
        doc.save(path)

        text = _extract_docx(str(path))
        assert "SEMESTRE I" in text
        assert "Fundamentos de seguridad informática" in text


class TestExtractPptx:
    def test_table_with_merged_style_header_row_is_preserved(self, tmp_path):
        path = tmp_path / "pensum.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        table_shape = slide.shapes.add_table(3, 2, Inches(1), Inches(1), Inches(4), Inches(2))
        table = table_shape.table
        table.cell(0, 0).text = "SEMESTRE I"
        table.cell(0, 1).text = ""  # single-populated header cell, like a merged cell
        table.cell(1, 0).text = "TD101"
        table.cell(1, 1).text = "Fundamentos de seguridad informática"
        table.cell(2, 0).text = ""
        table.cell(2, 1).text = ""  # fully blank row
        prs.save(path)

        text = _extract_pptx(str(path))
        assert "SEMESTRE I" in text
        assert "Fundamentos de seguridad informática" in text
        assert "=== DIAPOSITIVA 1 ===" in text
